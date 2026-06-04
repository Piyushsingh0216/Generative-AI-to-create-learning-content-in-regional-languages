from pathlib import Path
from typing import Set

import pdfplumber
import pytesseract
from docx import Document
from PIL import Image
from pptx import Presentation

from app.core.config import settings
from app.utils.exceptions import TextExtractionError, UnsupportedFileTypeError


SUPPORTED_EXTENSIONS: Set[str] = {".pdf", ".docx", ".pptx", ".txt", ".jpg", ".jpeg", ".png"}

if settings.tesseract_cmd:
    pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd


def _extract_pdf(file_path: Path) -> str:
    pages = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if not text or len(text.strip()) < 50:  # If little text, try OCR
                try:
                    # Use OCR on the page
                    img = page.to_image(resolution=300).original  # PIL Image
                    text = pytesseract.image_to_string(img)
                except Exception:
                    pass  # Keep original text
            pages.append(text or "")
    return "\n".join(pages)


def _extract_docx(file_path: Path) -> str:
    document = Document(file_path)
    return "\n".join(para.text for para in document.paragraphs if para.text)


def _extract_pptx(file_path: Path) -> str:
    presentation = Presentation(file_path)
    chunks = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
    return "\n".join(chunks)


def _extract_txt(file_path: Path) -> str:
    return file_path.read_text(encoding="utf-8", errors="ignore")


def _extract_image(file_path: Path) -> str:
    image = Image.open(file_path)
    return pytesseract.image_to_string(image)


def extract_text(file_path: Path) -> str:
    extension = file_path.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise UnsupportedFileTypeError(f"File type {extension} is not supported.")

    try:
        if extension == ".pdf":
            content = _extract_pdf(file_path)
        elif extension == ".docx":
            content = _extract_docx(file_path)
        elif extension == ".pptx":
            content = _extract_pptx(file_path)
        elif extension == ".txt":
            content = _extract_txt(file_path)
        else:
            content = _extract_image(file_path)
    except Exception as exc:
        raise TextExtractionError(f"Failed to extract text from {file_path.name}.") from exc

    normalized = "\n".join(line.strip() for line in content.splitlines() if line.strip())
    if not normalized:
        return "No extractable text found in the uploaded file."
    return normalized
